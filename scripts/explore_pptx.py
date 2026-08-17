#!/usr/bin/env python3
"""Explore PPTX structure for company boundaries."""
import json
import os
import re
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = os.path.join(os.path.expanduser("~"), "Desktop", "شروط صرف التعاقدات اغسطس 2026.pptx")
OUT = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "pptx_slides.json")


def shape_links(shape):
    links = []
    if not shape.has_text_frame:
        return links
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.hyperlink and run.hyperlink.address:
                links.append({"text": run.text.strip(), "url": run.hyperlink.address})
    return links


def shape_text(shape):
    if not shape.has_text_frame:
        return ""
    return shape.text.strip()


def extract():
    prs = Presentation(PPTX)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        links = []
        images = 0
        for shape in slide.shapes:
            t = shape_text(shape)
            if t:
                texts.append(t)
            links.extend(shape_links(shape))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    t = shape_text(child)
                    if t:
                        texts.append(t)
                    links.extend(shape_links(child))
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images += 1

        full_text = "\n".join(texts)
        slides.append({
            "slide": idx,
            "images": images,
            "links": links,
            "text_preview": full_text[:500],
            "text": full_text,
        })

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    with open(OUT, "w", encoding="utf-8") as f:
        json.dump({"slide_count": len(slides), "slides": slides}, f, ensure_ascii=False, indent=2)

    for s in slides[:30]:
        preview = s["text_preview"].replace("\n", " | ")[:90]
        print(f"{s['slide']:3d} imgs={s['images']} links={len(s['links'])} | {preview}")
    print(f"... total {len(slides)} slides -> {OUT}")


if __name__ == "__main__":
    import sys
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    extract()
