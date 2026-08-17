#!/usr/bin/env python3
"""Extract individual form images from PPTX for ALL company slides."""
from __future__ import annotations

import json
import os
import re
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, os.path.dirname(__file__))
from company_ranges import PAGE_RANGES, company_pages

BASE = os.path.join(os.path.dirname(os.path.dirname(__file__)))
PPTX_JSON = os.path.join(BASE, "data", "pptx_slides.json")
PPTX = os.path.join(os.path.expanduser("~"), "Desktop", "شروط صرف التعاقدات اغسطس 2026.pptx")
ASSETS = os.path.join(BASE, "data", "assets", "companies")
RULES = os.path.join(BASE, "data", "rules.json")

FORM_WORDS = ("نموذج", "روشت", "e-form", "eform", "علاج مسجل", "كربون", "مزمن", "form")
SKIP_LINE = ("الرجوع", "photo", "Photo", "Hotline", "http", "click to add")


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text.replace("\x0b", "\n")


def is_formish(line: str) -> bool:
    low = line.lower()
    return any(w in low or w in line for w in FORM_WORDS)


def clean_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip())[:120]


def is_valid_form_label(line: str) -> bool:
    if not line or any(s in line for s in SKIP_LINE):
        return False
    if is_formish(line):
        return True
    if re.match(r"^\d+\s*[-.)]", line):
        return False
    return len(line) >= 15 and any(w in line.lower() for w in ("approval", "prior", "موافق", "icare", "yodawy", "يوداو"))


def extract_formish_lines(blocks: list[str]) -> list[str]:
    found: list[str] = []
    seen: set[str] = set()
    for block in blocks:
        for line in block.split("\n"):
            line = clean_line(line)
            if not is_valid_form_label(line):
                continue
            if line not in seen:
                seen.add(line)
                found.append(line)
    return found


def extract_numbered_lines(blocks: list[str]) -> list[str]:
    found: list[str] = []
    seen: set[str] = set()
    for block in blocks:
        text = block.replace("\x0b", "\n")
        # Lines starting with digit-
        for line in text.split("\n"):
            line = clean_line(line)
            if not is_valid_form_label(line):
                continue
            if re.match(r"^\d+\s*[-.)]", line):
                if line not in seen:
                    seen.add(line)
                    found.append(line)
        # Inline numbered segments in one paragraph
        for m in re.finditer(r"\d+\s*[-.)]\s*[^0-9\n]+", text):
            line = clean_line(m.group(0))
            if is_valid_form_label(line) and line not in seen:
                seen.add(line)
                found.append(line)
    return found


def caption_shapes(blocks_with_pos: list[dict]) -> list[dict]:
    caps = []
    for b in blocks_with_pos:
        line = clean_line(b["text"])
        if re.match(r"^\d+\s*[-.)]", line) and is_valid_form_label(line):
            caps.append({**b, "text": line})
    return caps


def form_tags_from_label(label: str) -> list[str]:
    n = label.lower()
    tags = [label]
    if "اصفر" in n or "أصفر" in label:
        tags.extend(["أصفر", "اصفر", "yellow", "كربون", "نموذج أصفر", "بالكربون"])
    if "ازرق" in n or "أزرق" in label:
        tags.extend(["أزرق", "ازرق", "blue", "globemed"])
    if "ابيض" in n or "أبيض" in label:
        tags.extend(["أبيض", "ابيض", "white"])
    if "وان هيلث" in n or "one health" in n or "onehealth" in n:
        tags.extend(["one health", "onehealth", "وان هيلث", "روشتة one health", "one"])
    if "خارج" in n:
        tags.extend(["خارجية", "خارج", "external", "out of network"])
    if "e-form" in n or "eform" in n or "يواد" in n or "yodawy" in n or "يوداو" in n:
        tags.extend(["e-form", "eform", "yodawy", "يوادوي", "يوداوي"])
    if "مزمن" in n:
        tags.append("مزمن")
    if "شهري" in n:
        tags.extend(["شهري", "chronic"])
    if "e-prescription" in n or "eprescription" in n:
        tags.extend(["e-prescription", "prescription", "روشتة", "إلكترون"])
    if "بدون روشت" in n:
        tags.append("بدون روشتة")
    if "نموذج" in label:
        tags.append("نموذج")
    if "روشت" in label:
        tags.append("روشتة")
    if "مسجل" in label:
        tags.extend(["مسجل", "برنامج", "علاج مسجل"])
    return list(dict.fromkeys(tags))


def pick_fallback_label(blocks: list[str]) -> str:
    for block in blocks:
        for line in block.split("\n"):
            line = clean_line(line)
            if len(line) < 8 or any(s in line for s in SKIP_LINE):
                continue
            if is_formish(line):
                return line
    for block in blocks:
        for line in block.split("\n"):
            line = clean_line(line)
            if len(line) >= 8 and not any(s in line for s in SKIP_LINE):
                return line
    return ""


def line_sort_key(line: str) -> tuple:
    m = re.match(r"^(\d+)", line)
    return (int(m.group(1)) if m else 99, line)


def match_pics_to_labels(pics: list[dict], labels: list[str], cap_shapes: list[dict]) -> list[tuple[dict, str]]:
    if not pics or not labels:
        return []

    # Separate caption shapes with x position
    if len(cap_shapes) >= 2 and len(pics) >= 2:
        used: set[str] = set()
        pairs = []
        for pic in sorted(pics, key=lambda p: -p["left"]):
            best, best_dist = None, float("inf")
            for cap in cap_shapes:
                if cap["text"] in used:
                    continue
                dist = abs(cap["center_x"] - pic["center_x"])
                if dist < best_dist:
                    best_dist, best = dist, cap
            if best and best_dist <= pic["width"] * 1.25:
                used.add(best["text"])
                pairs.append((pic, best["text"]))
        if pairs:
            return pairs

    # Number-ordered RTL zip
    if len(labels) >= 2:
        sorted_pics = sorted(pics, key=lambda p: -p["left"])
        sorted_labels = sorted(labels, key=line_sort_key)
        n = min(len(sorted_pics), len(sorted_labels))
        return [(sorted_pics[i], sorted_labels[i]) for i in range(n)]

    # Single picture
    if len(pics) == 1:
        label = labels[0] if labels else ""
        return [(pics[0], label)]

    return []


def extract_slide_forms(slide, slide_num: int) -> list[dict]:
    pics = []
    text_blocks: list[str] = []
    text_shapes: list[dict] = []

    for shape in iter_shapes(slide.shapes):
        raw = shape_text(shape)
        if raw.strip():
            text_blocks.append(raw)
            text_shapes.append({
                "text": raw,
                "left": shape.left,
                "center_x": shape.left + shape.width / 2,
            })
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.height >= 700_000 and shape.width >= 700_000:
                pics.append({
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "center_x": shape.left + shape.width / 2,
                    "blob": shape.image.blob,
                    "ext": shape.image.ext or "png",
                })

    if not pics:
        return []

    labels = extract_numbered_lines(text_blocks)
    if not labels:
        labels = extract_formish_lines(text_blocks)
    cap_shapes = caption_shapes(text_shapes)
    pairs = match_pics_to_labels(pics, labels, cap_shapes)

    if not pairs and len(pics) == 1:
        fallback = pick_fallback_label(text_blocks)
        if fallback:
            pairs = [(pics[0], fallback)]

    # Multi-pic slide: unmatched pics get generic labels
    if len(pics) >= 2 and len(pairs) < len(pics):
        used_pics = {id(p[0]) for p in pairs}
        extra_labels = [l for l in labels if l not in {p[1] for p in pairs}]
        idx = 0
        for pic in sorted(pics, key=lambda p: -p["left"]):
            if id(pic) in used_pics:
                continue
            label = extra_labels[idx] if idx < len(extra_labels) else f"مستند — صفحة {slide_num}"
            idx += 1
            pairs.append((pic, label))

    forms = []
    for pic, label in pairs:
        if not label:
            label = pick_fallback_label(text_blocks) or f"مستند — صفحة {slide_num}"
        forms.append({
            "label": label[:120],
            "tags": form_tags_from_label(label),
            "blob": pic["blob"],
            "ext": pic["ext"],
            "slide": slide_num,
        })
    return forms


def form_mentions_text(form_name: str, text: str) -> bool:
    text_n = clean_line(text).lower()
    form_n = clean_line(form_name).lower()
    if form_n in text_n or text_n in form_n:
        return True
    for tag in form_tags_from_label(form_name):
        t = tag.lower()
        if len(t) >= 4 and t in text_n:
            return True
    keys = [form_n]
    if "أزرق" in form_name or "ازرق" in form_n:
        keys.append("ازرق")
    if "أصفر" in form_name or "اصفر" in form_n:
        keys.append("اصفر")
    if "روشت" in form_n:
        keys.append("روشت")
    if "yodawy" in form_n or "يوداو" in form_name or "يواد" in form_n:
        keys.extend(["yodawy", "يوداو", "يواد"])
    if "مسجل" in form_n:
        keys.append("مسجل")
    if "خارج" in form_n:
        keys.append("خارج")
    if "مزمن" in form_n or "شهري" in form_n:
        keys.extend(["مزمن", "شهري", "chronic"])
    if "موافق" in form_n or "prior" in form_n or "call center" in form_n:
        keys.extend(["موافق", "approval", "prior"])
    if "e-prescription" in form_n or "eprescription" in form_n:
        keys.extend(["e-prescription", "prescription", "روشتة"])
    if "sehaone" in form_n or "seha one" in form_n:
        keys.append("sehaone")
    if "petroshad" in form_n:
        keys.append("petroshad")
    if "egycare" in form_n:
        keys.append("egycare")
    if "bupa" in form_n:
        keys.append("bupa")
    if "جروب" in form_name or "group" in form_n:
        keys.extend(["جروب", "group", "company"])
    if "صرف مباشر" in form_name:
        keys.extend(["صرف", "مباشر", "direct"])
    return any(k in text_n for k in keys if len(k) >= 3)


def enrich_forms_from_company_list(company: dict, slide_texts: dict[int, str]) -> None:
    cid = company["id"]
    if cid not in PAGE_RANGES:
        return
    forms_list = company.get("forms") or []
    if not forms_list:
        return

    form_media = [m for m in (company.get("media") or []) if m.get("type") == "form"]
    if not form_media:
        return

    company_names = {
        n.lower()
        for n in (company.get("nameEn"), company.get("nameAr"), company.get("id"))
        if n
    }

    for form_name in forms_list:
        tags = list(dict.fromkeys([form_name, *form_tags_from_label(form_name)]))
        pages = [
            p for p in company_pages(cid)
            if form_mentions_text(form_name, slide_texts.get(p, ""))
        ]
        form_lower = form_name.lower()
        if not pages and any(n in form_lower for n in company_names):
            pages = sorted({m.get("page", 0) for m in form_media})
        if not pages and len(forms_list) == 1:
            pages = sorted({m.get("page", 0) for m in form_media})
        if not pages:
            continue
        targets = [m for m in form_media if m.get("page") in pages]
        for item in targets:
            item["formTags"] = list(dict.fromkeys([*(item.get("formTags") or []), *tags]))
            if len(forms_list) == 1 and (item.get("title") or "").startswith("مستند —"):
                item["title"] = form_name

    assign_generic_form_labels(company)


def assign_generic_form_labels(company: dict) -> None:
    forms_list = company.get("forms") or []
    if not forms_list:
        return

    form_media = [m for m in (company.get("media") or []) if m.get("type") == "form"]
    if not form_media:
        return

    titled = {
        m.get("title")
        for m in form_media
        if m.get("title") and not str(m.get("title")).startswith("مستند —")
    }
    remaining_forms = [f for f in forms_list if f not in titled]
    if not remaining_forms:
        return

    generic_items = sorted(
        [m for m in form_media if str(m.get("title", "")).startswith("مستند —")],
        key=lambda m: (m.get("page", 0), m.get("id", "")),
    )
    for item, form_name in zip(generic_items, remaining_forms):
        tags = list(dict.fromkeys([form_name, *form_tags_from_label(form_name)]))
        item["formTags"] = list(dict.fromkeys([*(item.get("formTags") or []), *tags]))
        item["title"] = form_name


def main():
    prs = Presentation(PPTX)
    with open(RULES, encoding="utf-8") as f:
        rules = json.load(f)

    slide_texts: dict[int, str] = {}
    if os.path.exists(PPTX_JSON):
        with open(PPTX_JSON, encoding="utf-8") as f:
            slide_texts = {s["slide"]: s.get("text", "") for s in json.load(f).get("slides", [])}

    form_media_by_company: dict[str, list[dict]] = {cid: [] for cid in PAGE_RANGES}
    slide_counts: dict[str, int] = {}

    for cid, (start, end) in PAGE_RANGES.items():
        folder = os.path.join(ASSETS, cid)
        os.makedirs(folder, exist_ok=True)
        for slide_num in range(start, end + 1):
            extracted = extract_slide_forms(prs.slides[slide_num - 1], slide_num)
            if not extracted:
                continue
            slide_counts[cid] = slide_counts.get(cid, 0) + 1
            for idx, item in enumerate(extracted, 1):
                fname = f"slide{slide_num:03d}_form{idx}.{item['ext']}"
                with open(os.path.join(folder, fname), "wb") as out:
                    out.write(item["blob"])
                form_media_by_company[cid].append({
                    "id": f"{cid}-s{slide_num}-form{idx}",
                    "type": "form",
                    "title": item["label"],
                    "formTags": item["tags"],
                    "url": f"/assets/companies/{cid}/{fname}",
                    "page": slide_num,
                    "links": [],
                })

    for company in rules["companies"]:
        cid = company["id"]
        new_forms = form_media_by_company.get(cid, [])
        if not new_forms:
            continue
        form_slides = {m["page"] for m in new_forms}
        kept = [
            m for m in (company.get("media") or [])
            if m.get("type") not in ("form", "photo", "card") or m.get("page") not in form_slides
        ]
        merged = kept + new_forms
        merged.sort(key=lambda m: (m.get("page", 0), 0 if m.get("type") == "form" else 1, m.get("id", "")))
        company["media"] = merged
        enrich_forms_from_company_list(company, slide_texts)

    with open(RULES, "w", encoding="utf-8") as f:
        json.dump(rules, f, ensure_ascii=False, indent=2)

    total = sum(len(v) for v in form_media_by_company.values())
    print(f"Extracted {total} form images across {sum(1 for v in form_media_by_company.values() if v)} companies")
    for cid in PAGE_RANGES:
        n = len(form_media_by_company.get(cid, []))
        if n:
            print(f"  {cid}: {n} forms on {slide_counts.get(cid,0)} slides")


if __name__ == "__main__":
    main()
